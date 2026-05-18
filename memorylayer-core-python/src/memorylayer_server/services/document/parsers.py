"""Content parsers for document ingestion.

Each parser converts raw bytes into a list of ContentChunk objects.
Optional dependencies are imported with graceful fallbacks — a missing
library disables the parser rather than crashing at import time.
"""

import logging
import os
from abc import ABC, abstractmethod

from pydantic import BaseModel

from ...models.document import DocumentExtractionOptions, DocumentType

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Optional dependency guards
# ---------------------------------------------------------------------------

try:
    import pymupdf  # type: ignore[import-untyped]
except ImportError:
    pymupdf = None  # type: ignore[assignment]
    logger.debug("pymupdf not available — PDF text extraction disabled")

try:
    from docx import Document as DocxDocument  # type: ignore[import-untyped]
except ImportError:
    DocxDocument = None  # type: ignore[assignment]
    logger.debug("python-docx not available — DOCX parsing disabled")

try:
    import mammoth  # type: ignore[import-untyped]
except ImportError:
    mammoth = None  # type: ignore[assignment]
    logger.debug("mammoth not available — DOCX-to-HTML fallback disabled")

try:
    from pptx import Presentation as PptxPresentation  # type: ignore[import-untyped]
except ImportError:
    PptxPresentation = None  # type: ignore[assignment]
    logger.debug("python-pptx not available — PPTX parsing disabled")

try:
    import html2text  # type: ignore[import-untyped]
except ImportError:
    html2text = None  # type: ignore[assignment]
    logger.debug("html2text not available — HTML parsing will use basic stripping")

# ---------------------------------------------------------------------------
# Domain model
# ---------------------------------------------------------------------------

# Extension-to-language map used by CodeParser
_EXTENSION_LANGUAGE_MAP: dict[str, str] = {
    ".py": "python",
    ".js": "javascript",
    ".ts": "typescript",
    ".jsx": "javascript",
    ".tsx": "typescript",
    ".java": "java",
    ".kt": "kotlin",
    ".go": "go",
    ".rs": "rust",
    ".c": "c",
    ".cpp": "cpp",
    ".cc": "cpp",
    ".h": "c",
    ".hpp": "cpp",
    ".cs": "csharp",
    ".rb": "ruby",
    ".php": "php",
    ".swift": "swift",
    ".sh": "bash",
    ".bash": "bash",
    ".zsh": "zsh",
    ".sql": "sql",
    ".r": "r",
    ".scala": "scala",
    ".lua": "lua",
    ".dart": "dart",
    ".yml": "yaml",
    ".yaml": "yaml",
    ".toml": "toml",
    ".json": "json",
    ".xml": "xml",
    ".html": "html",
    ".css": "css",
    ".scss": "scss",
    ".tf": "terraform",
}


class ContentChunk(BaseModel):
    """A parsed chunk of content from a document."""

    text: str
    page_number: int | None = None
    metadata: dict = {}


# ---------------------------------------------------------------------------
# Abstract base
# ---------------------------------------------------------------------------


class ContentParser(ABC):
    """Abstract base for document content parsers."""

    @abstractmethod
    async def parse(
        self,
        content: bytes,
        filename: str,
        options: DocumentExtractionOptions,
    ) -> list[ContentChunk]:
        """Parse raw bytes into a list of content chunks.

        Args:
            content: Raw file bytes
            filename: Original filename (used for type/language detection)
            options: Extraction options controlling chunking behaviour

        Returns:
            List of ContentChunk objects
        """
        pass


# ---------------------------------------------------------------------------
# Concrete parsers
# ---------------------------------------------------------------------------


def _split_text_into_chunks(text: str, chunk_size: int, chunk_overlap: int) -> list[str]:
    """Split text into overlapping fixed-size chunks.

    Splits first on paragraph boundaries, then on raw character limits.
    """
    if not text.strip():
        return []

    if len(text) <= chunk_size:
        return [text]

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks: list[str] = []
    current_parts: list[str] = []
    current_len = 0

    for para in paragraphs:
        para_len = len(para)
        if current_len + para_len > chunk_size and current_parts:
            chunks.append("\n\n".join(current_parts))
            # Keep overlap by retaining recent paragraphs
            overlap_parts: list[str] = []
            overlap_len = 0
            for p in reversed(current_parts):
                if overlap_len + len(p) > chunk_overlap:
                    break
                overlap_parts.insert(0, p)
                overlap_len += len(p)
            current_parts = overlap_parts
            current_len = overlap_len

        current_parts.append(para)
        current_len += para_len

    if current_parts:
        chunks.append("\n\n".join(current_parts))

    return chunks


class TextParser(ContentParser):
    """Parser for plain text and Markdown files."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        options: DocumentExtractionOptions,
    ) -> list[ContentChunk]:
        text = content.decode("utf-8", errors="replace")
        raw_chunks = _split_text_into_chunks(text, options.chunk_size, options.chunk_overlap)
        return [ContentChunk(text=chunk, page_number=i) for i, chunk in enumerate(raw_chunks)]


class CodeParser(ContentParser):
    """Parser for source code files.

    Preserves the entire file as a single chunk with language metadata
    so that the code structure is not fragmented across chunks.
    """

    async def parse(
        self,
        content: bytes,
        filename: str,
        options: DocumentExtractionOptions,
    ) -> list[ContentChunk]:
        text = content.decode("utf-8", errors="replace")
        ext = os.path.splitext(filename)[1].lower()
        language = _EXTENSION_LANGUAGE_MAP.get(ext, "text")

        if not text.strip():
            return []

        # Wrap in a fenced code block so downstream LLMs see syntax highlighting
        formatted = f"```{language}\n{text}\n```"
        return [
            ContentChunk(
                text=formatted,
                page_number=0,
                metadata={"language": language, "filename": filename},
            )
        ]


class HTMLParser(ContentParser):
    """Parser for HTML files using html2text for Markdown conversion."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        options: DocumentExtractionOptions,
    ) -> list[ContentChunk]:
        raw_html = content.decode("utf-8", errors="replace")

        if html2text is not None:
            converter = html2text.HTML2Text()
            converter.ignore_links = False
            converter.ignore_images = True
            text = converter.handle(raw_html)
        else:
            # Basic tag stripping fallback
            import re

            text = re.sub(r"<[^>]+>", " ", raw_html)
            text = re.sub(r"\s+", " ", text).strip()

        raw_chunks = _split_text_into_chunks(text, options.chunk_size, options.chunk_overlap)
        return [ContentChunk(text=chunk, page_number=i) for i, chunk in enumerate(raw_chunks)]


class PDFParser(ContentParser):
    """Parser for PDF files.

    Uses pymupdf to extract page text. If text extraction yields no content
    for a page (e.g., scanned image), the page is skipped (vision transcription
    requires an LLM vision service and is handled by the document processing task).
    """

    async def parse(
        self,
        content: bytes,
        filename: str,
        options: DocumentExtractionOptions,
    ) -> list[ContentChunk]:
        if pymupdf is None:
            logger.warning(
                "pymupdf not installed — cannot parse PDF '%s'. Install pymupdf to enable PDF support.",
                filename,
            )
            return [
                ContentChunk(
                    text=f"[PDF parsing unavailable: install pymupdf to process '{filename}']",
                    page_number=0,
                    metadata={"parse_error": "pymupdf_not_installed"},
                )
            ]

        chunks: list[ContentChunk] = []
        try:
            doc = pymupdf.open(stream=content, filetype="pdf")
            for page_no, page in enumerate(doc):
                text = page.get_text().strip()
                if not text:
                    logger.debug("PDF page %d of '%s' has no extractable text, skipping", page_no, filename)
                    continue
                page_chunks = _split_text_into_chunks(text, options.chunk_size, options.chunk_overlap)
                for chunk in page_chunks:
                    chunks.append(ContentChunk(text=chunk, page_number=page_no, metadata={"source_page": page_no}))
            doc.close()
        except Exception as e:
            logger.error("Failed to parse PDF '%s': %s", filename, e)
            return [
                ContentChunk(
                    text=f"[Failed to parse PDF '{filename}': {e}]",
                    page_number=0,
                    metadata={"parse_error": str(e)},
                )
            ]

        return chunks


class ImageParser(ContentParser):
    """Parser for image files.

    Image transcription requires an LLM with vision capability. Without one,
    this parser returns a placeholder chunk so the document record is still
    created and the image stored for future reprocessing.
    """

    async def parse(
        self,
        content: bytes,
        filename: str,
        options: DocumentExtractionOptions,
    ) -> list[ContentChunk]:
        # Transcription is performed during the background task phase when an
        # LLM vision service is available. Return a placeholder here.
        logger.debug("ImageParser: transcription deferred to task phase for '%s'", filename)
        return [
            ContentChunk(
                text=f"[Image file '{filename}' — transcription pending LLM vision processing]",
                page_number=0,
                metadata={"pending_vision": True, "filename": filename},
            )
        ]


class DocxParser(ContentParser):
    """Parser for DOCX files using python-docx or mammoth fallback."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        options: DocumentExtractionOptions,
    ) -> list[ContentChunk]:
        import io

        text = ""

        if DocxDocument is not None:
            try:
                doc = DocxDocument(io.BytesIO(content))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                text = "\n\n".join(paragraphs)
            except Exception as e:
                logger.warning("python-docx failed for '%s': %s", filename, e)

        if not text and mammoth is not None:
            try:
                result = mammoth.convert_to_markdown(io.BytesIO(content))
                text = result.value
            except Exception as e:
                logger.warning("mammoth fallback failed for '%s': %s", filename, e)

        if not text:
            logger.warning(
                "No DOCX parser available for '%s'. Install python-docx or mammoth to enable DOCX support.",
                filename,
            )
            return [
                ContentChunk(
                    text=f"[DOCX parsing unavailable: install python-docx or mammoth to process '{filename}']",
                    page_number=0,
                    metadata={"parse_error": "no_docx_parser"},
                )
            ]

        raw_chunks = _split_text_into_chunks(text, options.chunk_size, options.chunk_overlap)
        return [ContentChunk(text=chunk, page_number=i) for i, chunk in enumerate(raw_chunks)]


class PptxParser(ContentParser):
    """Parser for PPTX files using python-pptx."""

    async def parse(
        self,
        content: bytes,
        filename: str,
        options: DocumentExtractionOptions,
    ) -> list[ContentChunk]:
        if PptxPresentation is None:
            logger.warning(
                "python-pptx not installed — cannot parse PPTX '%s'. Install python-pptx to enable PPTX support.",
                filename,
            )
            return [
                ContentChunk(
                    text=f"[PPTX parsing unavailable: install python-pptx to process '{filename}']",
                    page_number=0,
                    metadata={"parse_error": "pptx_not_installed"},
                )
            ]

        import io

        chunks: list[ContentChunk] = []
        try:
            prs = PptxPresentation(io.BytesIO(content))
            for slide_no, slide in enumerate(prs.slides):
                parts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                if parts:
                    slide_text = "\n".join(parts)
                    chunks.append(
                        ContentChunk(
                            text=slide_text,
                            page_number=slide_no,
                            metadata={"slide": slide_no + 1},
                        )
                    )
        except Exception as e:
            logger.error("Failed to parse PPTX '%s': %s", filename, e)
            return [
                ContentChunk(
                    text=f"[Failed to parse PPTX '{filename}': {e}]",
                    page_number=0,
                    metadata={"parse_error": str(e)},
                )
            ]

        return chunks


# ---------------------------------------------------------------------------
# Parser registry
# ---------------------------------------------------------------------------

_PARSER_MAP: dict[DocumentType, ContentParser] = {
    DocumentType.TEXT: TextParser(),
    DocumentType.MARKDOWN: TextParser(),
    DocumentType.CODE: CodeParser(),
    DocumentType.HTML: HTMLParser(),
    DocumentType.PDF: PDFParser(),
    DocumentType.IMAGE: ImageParser(),
    DocumentType.DOCX: DocxParser(),
    DocumentType.PPTX: PptxParser(),
}

_FALLBACK_PARSER = TextParser()


def get_parser(document_type: DocumentType) -> ContentParser:
    """Return the appropriate parser for a given document type.

    Falls back to TextParser for unknown types.

    Args:
        document_type: The document type to look up

    Returns:
        ContentParser instance for that type
    """
    return _PARSER_MAP.get(document_type, _FALLBACK_PARSER)


# ---------------------------------------------------------------------------
# Type detection helpers
# ---------------------------------------------------------------------------

_MIME_TO_TYPE: dict[str, DocumentType] = {
    "text/plain": DocumentType.TEXT,
    "text/markdown": DocumentType.MARKDOWN,
    "text/html": DocumentType.HTML,
    "application/pdf": DocumentType.PDF,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": DocumentType.DOCX,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": DocumentType.PPTX,
    "image/jpeg": DocumentType.IMAGE,
    "image/png": DocumentType.IMAGE,
    "image/gif": DocumentType.IMAGE,
    "image/webp": DocumentType.IMAGE,
    "image/bmp": DocumentType.IMAGE,
    "image/tiff": DocumentType.IMAGE,
}

_EXTENSION_TO_TYPE: dict[str, DocumentType] = {
    ".txt": DocumentType.TEXT,
    ".md": DocumentType.MARKDOWN,
    ".markdown": DocumentType.MARKDOWN,
    ".html": DocumentType.HTML,
    ".htm": DocumentType.HTML,
    ".pdf": DocumentType.PDF,
    ".docx": DocumentType.DOCX,
    ".pptx": DocumentType.PPTX,
    ".jpg": DocumentType.IMAGE,
    ".jpeg": DocumentType.IMAGE,
    ".png": DocumentType.IMAGE,
    ".gif": DocumentType.IMAGE,
    ".webp": DocumentType.IMAGE,
    ".bmp": DocumentType.IMAGE,
    ".tiff": DocumentType.IMAGE,
    ".tif": DocumentType.IMAGE,
}

# Code file extensions map to DocumentType.CODE
for _ext in _EXTENSION_LANGUAGE_MAP:
    if _ext not in _EXTENSION_TO_TYPE:
        _EXTENSION_TO_TYPE[_ext] = DocumentType.CODE


def detect_document_type(filename: str, mime_type: str | None = None) -> DocumentType:
    """Detect document type from filename extension or MIME type.

    MIME type takes priority over extension when both are available.

    Args:
        filename: Original filename
        mime_type: Optional MIME type from upload

    Returns:
        DocumentType enum value; defaults to TEXT if unrecognised
    """
    if mime_type:
        detected = _MIME_TO_TYPE.get(mime_type.lower().split(";")[0].strip())
        if detected:
            return detected

    ext = os.path.splitext(filename)[1].lower()
    return _EXTENSION_TO_TYPE.get(ext, DocumentType.TEXT)
